import tkinter as tk
from tkinter import scrolledtext, filedialog, messagebox, ttk, colorchooser, font
import threading
import os
import json
import random
import pandas as pd
import numpy as np
import matplotlib

matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.figure import Figure
from PIL import Image, ImageTk
from llama_cpp import Llama
import tempfile
from scipy import stats
from scipy.interpolate import interp1d
import sqlite3
from datetime import datetime
import re
import traceback
import warnings
from pptx import Presentation
from pptx.util import Inches
import shutil

warnings.filterwarnings('ignore')


# ========== ИНТЕГРИРОВАННЫЕ АЛГОРИТМЫ ИЗ well_performance_analysis ==========

class WellPlotAlgorithms:
    """Класс со всеми алгоритмами построения графиков из твоего модуля"""

    def __init__(self):
        self.column_names = {
            'date': 'Дата',
            'well': '№ скважины',
            'oil_rate': 'Дебит нефти за последний месяц, т/сут',
            'water_rate': 'Дебит воды за последний месяц, м3/сут',
            'watercut': 'Обводненность за посл.месяц, % (вес)',
            'bottomhole_pressure': 'Забойное давление, атм',
            'reservoir_pressure': 'Пластовое давление, атм'
        }
        self.calculation_params = {
            'reservoir_pressure': 250.0,
            'bubble_point_pressure': 150.0,
            'productivity_index': 2.5,
            'skin_factor': 0.0,
            'permeability': 100.0,
            'thickness': 10.0,
            'viscosity': 2.0,
            'volume_factor': 1.2,
            'drainage_radius': 300.0,
            'well_radius': 0.1,
            'density': 850.0
        }

    def auto_detect_columns(self, df):
        """Автоматическое определение колонок в DataFrame"""
        if df is None:
            return

        column_mapping = {
            'date': ['дата', 'date', 'время', 'time', 'месяц', 'год'],
            'well': ['скважин', 'well', '№', 'номер', 'name', 'скв'],
            'oil_rate': ['нефт', 'oil', 'дебит нефт', 'qн', 'q_oil', 'дебит'],
            'water_rate': ['вод', 'water', 'дебит вод', 'qв', 'q_water'],
            'watercut': ['обвод', 'watercut', 'влаг', 'содержан', 'water cut', '%'],
            'bottomhole_pressure': ['забойн', 'bottomhole', 'pзаб', 'p_заб', 'давл. забой'],
            'reservoir_pressure': ['пластов', 'reservoir', 'pпл', 'p_пл', 'давл. пласт']
        }

        detected_cols = {}
        for std_col, keywords in column_mapping.items():
            for col in df.columns:
                col_lower = str(col).lower()
                for keyword in keywords:
                    if keyword in col_lower:
                        detected_cols[std_col] = col
                        print(f"Автоопределение: {std_col} -> {col}")
                        break
                if std_col in detected_cols:
                    break

        for std_col, actual_col in detected_cols.items():
            self.column_names[std_col] = actual_col

        print(f"Итоговые имена колонок: {self.column_names}")

    def prepare_data(self, df):
        """Подготовка данных: преобразование типов"""
        if df is None:
            return df

        df = df.copy()
        date_col = self.column_names['date']
        if date_col in df.columns:
            df[date_col] = pd.to_datetime(df[date_col], errors='coerce')

        numeric_cols = [
            self.column_names.get('oil_rate'),
            self.column_names.get('water_rate'),
            self.column_names.get('watercut'),
            self.column_names.get('bottomhole_pressure'),
            self.column_names.get('reservoir_pressure')
        ]

        for col in numeric_cols:
            if col and col in df.columns:
                df[col] = pd.to_numeric(df[col], errors='coerce')

        return df

    def perform_chen_analysis(self, wor, d_log_wor_dt, time, oil_rate, watercut, well_name):
        """Выполнение детального анализа по методу Чена"""
        analysis_text = "АНАЛИЗ ПРИЧИН ОБВОДНЕНИЯ\n"
        analysis_text += "=" * 50 + "\n\n"
        analysis_text += f"Скважина: {well_name}\n"
        analysis_text += f"Период анализа: {len(time)} дней\n\n"

        # Базовые статистики
        analysis_text += "СТАТИСТИКА WOR:\n"
        analysis_text += "-" * 30 + "\n"
        analysis_text += f"Начальный WOR: {wor[0]:.3f} м³/т\n"
        analysis_text += f"Конечный WOR: {wor[-1]:.3f} м³/т\n"
        analysis_text += f"Рост WOR: в {wor[-1] / wor[0]:.1f} раз\n\n"

        # Анализ производной
        mean_derivative = np.mean(np.abs(d_log_wor_dt))
        max_derivative = np.max(np.abs(d_log_wor_dt))
        std_derivative = np.std(np.abs(d_log_wor_dt))

        analysis_text += "АНАЛИЗ ПРОИЗВОДНОЙ:\n"
        analysis_text += "-" * 30 + "\n"
        analysis_text += f"Средняя производная: {mean_derivative:.6f} 1/день\n"
        analysis_text += f"Максимальная производная: {max_derivative:.6f} 1/день\n"
        analysis_text += f"Стандартное отклонение: {std_derivative:.6f} 1/день\n\n"

        # Линейная регрессия для определения тренда
        if len(time) >= 3:
            slope, intercept, r_value, p_value, std_err = stats.linregress(time, np.log(wor))

            analysis_text += "АНАЛИЗ ТРЕНДА:\n"
            analysis_text += "-" * 30 + "\n"
            analysis_text += f"Наклон тренда: {slope:.6f} 1/день\n"
            analysis_text += f"Коэффициент детерминации R²: {r_value ** 2:.3f}\n"
            analysis_text += f"Время удвоения WOR: {np.log(2) / abs(slope):.1f} дней\n\n"

        # Диагностика по классификации Чена
        analysis_text += "ДИАГНОСТИКА ПО МЕТОДУ ЧЕНА:\n"
        analysis_text += "-" * 30 + "\n"

        diagnosis = ""
        cause = ""
        recommendations = []

        if mean_derivative < 0.001:
            diagnosis = "ТИП I: Стабильное поведение"
            cause = "Прорыв воды по поровому пространству\n(медленный фронт вытеснения)"
            recommendations = [
                "✓ Контроль режима работы скважины",
                "✓ Возможно увеличение отборов",
                "✓ Мониторинг пластового давления"
            ]
        elif 0.001 <= mean_derivative < 0.01:
            diagnosis = "ТИП II: Умеренный рост"
            cause = "Прорыв воды по высокопроницаемому пропластку"
            recommendations = [
                "✓ Провести ГИС для определения интервала прорыва",
                "✓ Рассмотреть возможность изоляции пропластка",
                "✓ Оптимизировать режим работы"
            ]
        elif 0.01 <= mean_derivative < 0.1:
            diagnosis = "ТИП III: Быстрый рост"
            cause = "Негерметичность эксплуатационной колонны\nили заколонная циркуляция"
            recommendations = [
                "✓ Проверить герметичность колонны",
                "✓ Выполнить ремонтные работы",
                "✓ Рассмотреть возможность цементирования"
            ]
        elif mean_derivative >= 0.1:
            diagnosis = "ТИП IV: Критический рост"
            cause = "Катастрофический прорыв воды\nили межпластовые перетоки"
            recommendations = [
                "✓ СРОЧНАЯ оптимизация режима работы",
                "✓ Рассмотреть временную остановку скважины",
                "✓ Планировать капитальный ремонт"
            ]

        analysis_text += f"Диагноз: {diagnosis}\n"
        analysis_text += f"\nВероятная причина:\n{cause}\n"

        analysis_text += "\nРЕКОМЕНДАЦИИ:\n"
        analysis_text += "-" * 30 + "\n"
        for rec in recommendations:
            analysis_text += f"{rec}\n"

        # Дополнительный анализ
        analysis_text += "\nДОПОЛНИТЕЛЬНЫЙ АНАЛИЗ:\n"
        analysis_text += "-" * 30 + "\n"

        # Анализ корреляции
        if len(oil_rate) == len(wor):
            oil_corr = np.corrcoef(oil_rate, wor)[0, 1]
            analysis_text += f"Корреляция нефть-WOR: {oil_corr:.3f}\n"

        # Анализ обводненности
        if watercut is not None and len(watercut) > 0:
            analysis_text += f"Начальная обводненность: {watercut[0]:.1f}%\n"
            analysis_text += f"Конечная обводненность: {watercut[-1]:.1f}%\n"
            analysis_text += f"Прирост обводненности: {watercut[-1] - watercut[0]:.1f}%\n"

        # Прогноз
        analysis_text += "\nПРОГНОЗ:\n"
        analysis_text += "-" * 30 + "\n"
        if 'slope' in locals():
            days_to_90_percent = (np.log(9) - np.log(wor[-1] / wor[0])) / slope if slope > 0 else "∞"
            analysis_text += f"До достижения WOR=9 (90% обвод.): {days_to_90_percent:.0f} дней\n"

        return analysis_text

    def plot_chen(self, df, well_name, output_path, smooth_days=30):
        """
        График Чена для диагностики обводнения
        Возвращает путь к сохраненному графику и текст анализа
        """
        try:
            # Фильтруем по скважине
            well_col = self.column_names['well']
            well_data = df[df[well_col] == well_name].sort_values(self.column_names['date'])

            print(f"Chen: найдено {len(well_data)} записей для {well_name}")

            if len(well_data) < 5:
                raise ValueError(
                    f"Недостаточно данных для скважины {well_name} (минимум 5 точек, есть {len(well_data)})")

            # Создаем фигуру с двумя подграфиками
            fig = plt.figure(figsize=(12, 8), dpi=100)

            # Подграфик 1: WOR в логарифмических координатах
            ax1 = fig.add_subplot(211)

            # Подграфик 2: Производная WOR
            ax2 = fig.add_subplot(212)

            # Получаем данные
            dates = well_data[self.column_names['date']].values
            oil_rate = well_data[self.column_names['oil_rate']].values
            water_rate = well_data[self.column_names['water_rate']].values
            watercut = well_data[self.column_names['watercut']].values if self.column_names[
                                                                              'watercut'] in well_data.columns else None

            # Рассчитываем время в днях от начала
            if len(dates) > 0:
                try:
                    if isinstance(dates[0], np.datetime64):
                        days_from_start = np.array([(pd.Timestamp(d) - pd.Timestamp(dates[0])).days for d in dates])
                    else:
                        days_from_start = np.array([(d - dates[0]).days for d in dates])
                except:
                    days_from_start = np.arange(len(dates))
            else:
                days_from_start = np.array([])

            # Рассчитываем WOR (Water-Oil Ratio)
            with np.errstate(divide='ignore', invalid='ignore'):
                wor = np.where(oil_rate > 0.1, water_rate / oil_rate, np.nan)

            # Сглаживание
            smooth_window = smooth_days
            if smooth_window > 0 and len(wor) > smooth_window:
                wor_series = pd.Series(wor)
                wor_smooth = wor_series.rolling(window=min(smooth_window, len(wor) // 2),
                                                center=True, min_periods=1).mean().values
            else:
                wor_smooth = wor

            # Рассчитываем производную
            valid_mask = ~np.isnan(wor_smooth)

            if np.sum(valid_mask) >= 3:
                log_wor = np.log(wor_smooth[valid_mask])
                time_valid = days_from_start[valid_mask]
                d_log_wor_dt = np.gradient(log_wor, time_valid)

                # График 1: WOR
                ax1.semilogy(time_valid, wor_smooth[valid_mask], 'b-', linewidth=2,
                             marker='o', markersize=4, label='WOR (сглаженный)')
                ax1.semilogy(days_from_start, wor, 'ro', markersize=3, alpha=0.3, label='WOR (факт)')

                # Тренд
                if len(time_valid) > 1:
                    slope, intercept, r_value, p_value, std_err = stats.linregress(time_valid, log_wor)
                    trend_line = np.exp(intercept + slope * time_valid)
                    ax1.semilogy(time_valid, trend_line, 'g--', linewidth=1.5, alpha=0.7,
                                 label=f'Тренд (наклон={slope:.4f})')

                ax1.set_xlabel('Время, дни', fontsize=10)
                ax1.set_ylabel('WOR (вода/нефть)', fontsize=10)
                ax1.set_title(f'График Чена: {well_name}', fontsize=12, fontweight='bold')
                ax1.grid(True, alpha=0.3, which='both')
                ax1.legend(loc='best')

                # График 2: Производная
                ax2.semilogy(time_valid[1:-1], np.abs(d_log_wor_dt[1:-1]), 'g-', linewidth=2,
                             marker='s', markersize=4, label='|d(ln WOR)/dt|')

                # Диагностические зоны
                ax2.axhspan(0, 0.001, alpha=0.2, color='green', label='Медленный рост')
                ax2.axhspan(0.001, 0.01, alpha=0.2, color='yellow', label='Средний рост')
                ax2.axhspan(0.01, 0.1, alpha=0.2, color='orange', label='Быстрый рост')
                ax2.axhspan(0.1, 10, alpha=0.2, color='red', label='Критический рост')

                ax2.set_xlabel('Время, дни', fontsize=10)
                ax2.set_ylabel('|d(ln WOR)/dt|, 1/день', fontsize=10)
                ax2.set_title('Производная WOR', fontsize=12, fontweight='bold')
                ax2.grid(True, alpha=0.3, which='both')
                ax2.legend(loc='best', fontsize=8)

                # Выполняем анализ
                analysis_text = self.perform_chen_analysis(
                    wor_smooth[valid_mask],
                    d_log_wor_dt,
                    time_valid,
                    oil_rate[valid_mask] if len(oil_rate[valid_mask]) == len(time_valid) else oil_rate[
                        :len(time_valid)],
                    watercut[:len(time_valid)] if watercut is not None else None,
                    well_name
                )
            else:
                ax1.plot(days_from_start, wor, 'bo-', markersize=5)
                ax1.set_title(f'График Чена (недостаточно данных)')
                analysis_text = "Недостаточно данных для анализа"

            plt.tight_layout()
            fig.savefig(output_path, dpi=100, bbox_inches='tight')
            plt.close(fig)
            return output_path, analysis_text

        except Exception as e:
            plt.close('all')
            raise e

    def plot_ipr(self, df, well_name, output_path, params=None):
        """Построение IPR диаграммы"""
        try:
            if params:
                for key, value in params.items():
                    self.calculation_params[key] = value

            well_col = self.column_names['well']
            well_data = df[df[well_col] == well_name].sort_values(self.column_names['date'])

            print(f"IPR: найдено {len(well_data)} записей для {well_name}")

            if len(well_data) < 1:
                raise ValueError(f"Нет данных для скважины {well_name}")

            Pe = self.calculation_params['reservoir_pressure']
            Pb = self.calculation_params['bubble_point_pressure']
            J = self.calculation_params['productivity_index']

            # Используем фактическое пластовое давление если есть
            if self.column_names['reservoir_pressure'] in well_data.columns:
                Pe_actual = well_data[self.column_names['reservoir_pressure']].iloc[-1]
                if not pd.isna(Pe_actual):
                    Pe = Pe_actual

            fig = Figure(figsize=(10, 6), dpi=100)
            ax = fig.add_subplot(111)

            Pwf_range = np.linspace(0, Pe, 50)

            # Линейная модель
            q_linear = J * (Pe - Pwf_range)

            # Модель Вогела
            q_vogel = []
            for Pwf in Pwf_range:
                if Pwf >= Pb:
                    q = J * (Pe - Pwf)
                else:
                    q = J * (Pe - Pb) + (J * Pb / 1.8) * (1 - 0.2 * (Pwf / Pb) - 0.8 * (Pwf / Pb) ** 2)
                q_vogel.append(q)

            # Строим обе модели
            ax.plot(q_linear, Pwf_range, 'b--', linewidth=2, label='Линейная модель', alpha=0.7)
            ax.plot(q_vogel, Pwf_range, 'r-', linewidth=3, label='Модель Вогела')

            # Фактические данные
            if self.column_names['bottomhole_pressure'] in well_data.columns:
                actual_pwf = well_data[self.column_names['bottomhole_pressure']].values
                actual_q = well_data[self.column_names['oil_rate']].values

                valid_mask = ~pd.isna(actual_pwf) & ~pd.isna(actual_q) & (actual_q > 0)
                if np.any(valid_mask):
                    ax.scatter(actual_q[valid_mask], actual_pwf[valid_mask],
                               c='red', s=50, marker='o', label='Фактические данные', zorder=5)

                    if np.sum(valid_mask) >= 2:
                        slope, intercept, r_value, p_value, std_err = stats.linregress(
                            actual_q[valid_mask], actual_pwf[valid_mask]
                        )
                        q_fit = np.linspace(0, max(actual_q[valid_mask]) * 1.2, 50)
                        pwf_fit = slope * q_fit + intercept
                        ax.plot(q_fit, pwf_fit, 'r--', alpha=0.5,
                                label=f'R²={r_value ** 2:.3f}')

            ax.set_xlabel('Дебит нефти, т/сут', fontsize=11)
            ax.set_ylabel('Забойное давление, атм', fontsize=11)
            ax.set_title(f'IPR Диаграмма: {well_name}', fontsize=13, fontweight='bold')
            ax.grid(True, alpha=0.3)
            ax.legend(loc='upper right')
            ax.invert_yaxis()
            ax.set_xlim(left=0)

            # Линии давления
            ax.axhline(y=Pb, color='orange', linestyle='--', alpha=0.7, label=f'Pнас = {Pb:.1f}')
            ax.axhline(y=Pe, color='green', linestyle='--', alpha=0.7, label=f'Pпл = {Pe:.1f}')

            fig.savefig(output_path, dpi=100, bbox_inches='tight')
            plt.close(fig)
            return output_path, "IPR диаграмма построена"

        except Exception as e:
            plt.close('all')
            raise e

    def plot_decline(self, df, well_name, output_path):
        """Кривая падения добычи"""
        try:
            well_col = self.column_names['well']
            well_data = df[df[well_col] == well_name].sort_values(self.column_names['date'])

            print(f"Decline: найдено {len(well_data)} записей для {well_name}")

            if len(well_data) < 3:
                raise ValueError(
                    f"Недостаточно данных для скважины {well_name} (минимум 3 точки, есть {len(well_data)})")

            fig = Figure(figsize=(10, 6), dpi=100)
            ax = fig.add_subplot(111)

            oil_rate = well_data[self.column_names['oil_rate']].values
            time_months = np.arange(len(oil_rate))

            ax.plot(time_months, oil_rate, 'bo-', linewidth=2, markersize=6, label='Фактические данные')

            # Простая экспоненциальная аппроксимация
            if len(oil_rate) > 3:
                try:
                    log_rate = np.log(oil_rate[oil_rate > 0])
                    time_valid = time_months[oil_rate > 0]
                    if len(log_rate) > 2:
                        slope, intercept, r_value, p_value, std_err = stats.linregress(time_valid, log_rate)
                        q_fit = np.exp(intercept + slope * time_months)
                        ax.plot(time_months, q_fit, 'r--', linewidth=2,
                                label=f'Эксп. тренд (Di={-slope:.3f})')
                except:
                    pass

            ax.set_xlabel('Время, месяцы', fontsize=10)
            ax.set_ylabel('Дебит нефти, т/сут', fontsize=10)
            ax.set_title(f'Decline Curve: {well_name}', fontsize=12, fontweight='bold')
            ax.grid(True, alpha=0.3)
            ax.legend()

            fig.savefig(output_path, dpi=100, bbox_inches='tight')
            plt.close(fig)
            return output_path, "Кривая падения построена"

        except Exception as e:
            plt.close('all')
            raise e

    def plot_standard(self, df, well_name, output_path, graph_type='oil_rate'):
        """Стандартные графики (дебит нефти или обводненность)"""
        try:
            well_col = self.column_names['well']
            well_data = df[df[well_col] == well_name].sort_values(self.column_names['date'])

            print(f"Standard: найдено {len(well_data)} записей для {well_name}")

            if len(well_data) < 1:
                raise ValueError(f"Нет данных для скважины {well_name}")

            fig = Figure(figsize=(10, 6), dpi=100)
            ax = fig.add_subplot(111)

            dates = well_data[self.column_names['date']]

            if graph_type == 'oil_rate':
                values = well_data[self.column_names['oil_rate']]
                ax.plot(dates, values, 'ro-', linewidth=2, markersize=5)
                ax.set_ylabel("Дебит нефти, т/сут", fontsize=10)
                ax.set_title(f"Дебит нефти: {well_name}", fontsize=12, fontweight='bold')
                analysis_text = f"Дебит нефти {well_name}: средний {np.mean(values):.1f} т/сут"
            elif graph_type == 'watercut':
                values = well_data[self.column_names['watercut']]
                ax.plot(dates, values, 'g^-', linewidth=2, markersize=5)
                ax.set_ylabel("Обводненность, %", fontsize=10)
                ax.set_title(f"Обводненность: {well_name}", fontsize=12, fontweight='bold')
                ax.axhline(y=90, color='r', linestyle='--', alpha=0.5, label='Критическая (90%)')
                ax.legend()
                analysis_text = f"Обводненность {well_name}: средняя {np.mean(values):.1f}%"

            ax.set_xlabel("Дата", fontsize=10)
            ax.grid(True, alpha=0.3)
            ax.tick_params(axis='x', rotation=45)
            fig.tight_layout()

            fig.savefig(output_path, dpi=100, bbox_inches='tight')
            plt.close(fig)
            return output_path, analysis_text

        except Exception as e:
            plt.close('all')
            raise e


# ========== КЛАСС ДЛЯ НАСТРОЙКИ ИНТЕРФЕЙСА ==========

class SettingsDialog:
    """Диалог настройки цветов и шрифтов интерфейса"""

    def __init__(self, parent, colors, font_settings, callback):
        self.parent = parent
        self.colors = colors.copy()
        self.font_settings = font_settings.copy()
        self.callback = callback

        self.dialog = tk.Toplevel(parent)
        self.dialog.title("Настройки интерфейса")
        self.dialog.geometry("500x600")
        self.dialog.configure(bg=colors['bg'])
        self.dialog.transient(parent)
        self.dialog.grab_set()

        self.create_widgets()

    def create_widgets(self):
        # Заголовок
        title_label = tk.Label(
            self.dialog,
            text="🦀 Настройка внешнего вида",
            font=("Arial", 14, "bold"),
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        title_label.pack(pady=20)

        # ===== НАСТРОЙКА ЦВЕТОВ =====
        colors_frame = tk.LabelFrame(
            self.dialog,
            text="Цвета интерфейса",
            bg=self.colors['bg'],
            fg=self.colors['text'],
            font=("Arial", 11, "bold")
        )
        colors_frame.pack(fill='x', padx=20, pady=10)

        # Цвет фона
        bg_frame = tk.Frame(colors_frame, bg=self.colors['bg'])
        bg_frame.pack(fill='x', pady=5)

        tk.Label(bg_frame, text="Фон:", bg=self.colors['bg'], fg=self.colors['text']).pack(side='left', padx=10)
        self.bg_color_btn = tk.Button(
            bg_frame,
            text="Выбрать цвет",
            command=lambda: self.choose_color('bg'),
            bg=self.colors['bg'],
            fg='white'
        )
        self.bg_color_btn.pack(side='left', padx=10)
        self.bg_color_label = tk.Label(
            bg_frame,
            text=self.colors['bg'],
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        self.bg_color_label.pack(side='left')

        # Цвет панелей
        panel_frame = tk.Frame(colors_frame, bg=self.colors['bg'])
        panel_frame.pack(fill='x', pady=5)

        tk.Label(panel_frame, text="Панели:", bg=self.colors['bg'], fg=self.colors['text']).pack(side='left', padx=10)
        self.panel_color_btn = tk.Button(
            panel_frame,
            text="Выбрать цвет",
            command=lambda: self.choose_color('panel'),
            bg=self.colors['panel'],
            fg='white'
        )
        self.panel_color_btn.pack(side='left', padx=10)
        self.panel_color_label = tk.Label(
            panel_frame,
            text=self.colors['panel'],
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        self.panel_color_label.pack(side='left')

        # Цвет текста
        text_frame = tk.Frame(colors_frame, bg=self.colors['bg'])
        text_frame.pack(fill='x', pady=5)

        tk.Label(text_frame, text="Текст:", bg=self.colors['bg'], fg=self.colors['text']).pack(side='left', padx=10)
        self.text_color_btn = tk.Button(
            text_frame,
            text="Выбрать цвет",
            command=lambda: self.choose_color('text'),
            bg=self.colors['text'],
            fg='black'
        )
        self.text_color_btn.pack(side='left', padx=10)
        self.text_color_label = tk.Label(
            text_frame,
            text=self.colors['text'],
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        self.text_color_label.pack(side='left')

        # Цвет акцента (синий)
        accent_frame = tk.Frame(colors_frame, bg=self.colors['bg'])
        accent_frame.pack(fill='x', pady=5)

        tk.Label(accent_frame, text="Акцент:", bg=self.colors['bg'], fg=self.colors['text']).pack(side='left', padx=10)
        self.accent_color_btn = tk.Button(
            accent_frame,
            text="Выбрать цвет",
            command=lambda: self.choose_color('accent_blue'),
            bg=self.colors['accent_blue'],
            fg='white'
        )
        self.accent_color_btn.pack(side='left', padx=10)
        self.accent_color_label = tk.Label(
            accent_frame,
            text=self.colors['accent_blue'],
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        self.accent_color_label.pack(side='left')

        # ===== НАСТРОЙКА ШРИФТОВ =====
        font_frame = tk.LabelFrame(
            self.dialog,
            text="Шрифты",
            bg=self.colors['bg'],
            fg=self.colors['text'],
            font=("Arial", 11, "bold")
        )
        font_frame.pack(fill='x', padx=20, pady=10)

        # Семейство шрифта
        family_frame = tk.Frame(font_frame, bg=self.colors['bg'])
        family_frame.pack(fill='x', pady=5)

        tk.Label(family_frame, text="Шрифт:", bg=self.colors['bg'], fg=self.colors['text']).pack(side='left', padx=10)

        # Получаем список доступных шрифтов
        available_fonts = sorted(list(tk.font.families()))

        self.font_var = tk.StringVar(value=self.font_settings['family'])
        font_menu = ttk.Combobox(
            family_frame,
            textvariable=self.font_var,
            values=available_fonts[:50],  # Ограничим для скорости
            width=20
        )
        font_menu.pack(side='left', padx=10)

        # Размер шрифта
        size_frame = tk.Frame(font_frame, bg=self.colors['bg'])
        size_frame.pack(fill='x', pady=5)

        tk.Label(size_frame, text="Размер:", bg=self.colors['bg'], fg=self.colors['text']).pack(side='left', padx=10)

        self.size_var = tk.IntVar(value=self.font_settings['size'])
        size_spin = tk.Spinbox(
            size_frame,
            from_=8,
            to=24,
            textvariable=self.size_var,
            width=10
        )
        size_spin.pack(side='left', padx=10)

        # ===== ПРЕДПРОСМОТР =====
        preview_frame = tk.LabelFrame(
            self.dialog,
            text="Предпросмотр",
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        preview_frame.pack(fill='both', expand=True, padx=20, pady=10)

        self.preview_text = tk.Text(
            preview_frame,
            height=5,
            font=(self.font_var.get(), self.size_var.get()),
            bg=self.colors['bg'],
            fg=self.colors['text'],
            wrap=tk.WORD
        )
        self.preview_text.pack(fill='both', expand=True, padx=10, pady=10)
        self.preview_text.insert(1.0, "Пример текста в чате\nСкважина 109: дебит нефти 45.3 т/сут")
        self.preview_text.config(state='disabled')

        # Кнопки
        buttons_frame = tk.Frame(self.dialog, bg=self.colors['bg'])
        buttons_frame.pack(fill='x', padx=20, pady=20)

        tk.Button(
            buttons_frame,
            text="Применить",
            command=self.apply_settings,
            bg=self.colors['accent_green'],
            fg='white',
            font=("Arial", 11),
            padx=20
        ).pack(side='left', padx=10)

        tk.Button(
            buttons_frame,
            text="Отмена",
            command=self.dialog.destroy,
            bg=self.colors['accent_red'],
            fg='white',
            font=("Arial", 11),
            padx=20
        ).pack(side='right', padx=10)

    def choose_color(self, key):
        color = colorchooser.askcolor(title=f"Выберите цвет для {key}")[1]
        if color:
            self.colors[key] = color
            # Обновляем кнопку и метку
            getattr(self, f'{key}_color_btn').config(bg=color)
            getattr(self, f'{key}_color_label').config(text=color)
            # Обновляем предпросмотр
            self.update_preview()

    def update_preview(self):
        self.preview_text.config(
            bg=self.colors['bg'],
            fg=self.colors['text'],
            font=(self.font_var.get(), self.size_var.get())
        )

    def apply_settings(self):
        self.font_settings['family'] = self.font_var.get()
        self.font_settings['size'] = self.size_var.get()
        self.callback(self.colors, self.font_settings)
        self.dialog.destroy()


# ========== ДИАЛОГ ВЫБОРА МОДЕЛИ ==========

class ModelSelectionDialog:
    """Диалог выбора и загрузки модели GGUF"""

    def __init__(self, parent, current_path, callback):
        self.parent = parent
        self.current_path = current_path
        self.callback = callback
        self.model_path = current_path

        self.dialog = tk.Toplevel(parent)
        self.dialog.title("Выбор модели")
        self.dialog.geometry("500x300")
        self.dialog.transient(parent)
        self.dialog.grab_set()

        self.create_widgets()

    def create_widgets(self):
        # Заголовок
        title_label = tk.Label(
            self.dialog,
            text="🤖 Выберите файл модели",
            font=("Arial", 14, "bold")
        )
        title_label.pack(pady=20)

        # Текущий путь
        current_frame = tk.Frame(self.dialog)
        current_frame.pack(fill='x', padx=20, pady=10)

        tk.Label(current_frame, text="Текущая модель:", font=("Arial", 10)).pack(anchor='w')

        self.current_label = tk.Label(
            current_frame,
            text=self.current_path if self.current_path else "не выбрана",
            font=("Arial", 9),
            fg="blue",
            wraplength=450
        )
        self.current_label.pack(anchor='w', pady=5)

        # Кнопка выбора файла
        select_btn = tk.Button(
            self.dialog,
            text="📁 Выбрать файл GGUF",
            command=self.select_file,
            font=("Arial", 11),
            bg="lightblue",
            padx=20,
            pady=10
        )
        select_btn.pack(pady=10)

        # Информация
        info_text = "Файлы модели имеют расширение .gguf\n"
        info_text += "Рекомендуемый размер: 4-8 GB"
        info_label = tk.Label(
            self.dialog,
            text=info_text,
            font=("Arial", 9),
            fg="gray"
        )
        info_label.pack(pady=10)

        # Кнопки
        buttons_frame = tk.Frame(self.dialog)
        buttons_frame.pack(fill='x', padx=20, pady=20)

        tk.Button(
            buttons_frame,
            text="Сохранить и перезагрузить",
            command=self.apply,
            bg="lightgreen",
            font=("Arial", 10),
            padx=20
        ).pack(side='left', padx=10)

        tk.Button(
            buttons_frame,
            text="Отмена",
            command=self.dialog.destroy,
            bg="lightcoral",
            font=("Arial", 10),
            padx=20
        ).pack(side='right', padx=10)

    def select_file(self):
        filename = filedialog.askopenfilename(
            title="Выберите файл модели GGUF",
            filetypes=[("GGUF files", "*.gguf"), ("All files", "*.*")]
        )
        if filename:
            self.model_path = filename
            self.current_label.config(text=filename, fg="green")

    def apply(self):
        if self.model_path and os.path.exists(self.model_path):
            self.callback(self.model_path)
            self.dialog.destroy()
        else:
            messagebox.showerror("Ошибка", "Файл модели не найден!")


# ========== ОСНОВНОЙ КЛАСС АССИСТЕНТА ==========

class CrabAssistant:
    def __init__(self):
        self.window = tk.Tk()
        self.window.title("🦀 Нефтяной ассистент")
        self.window.geometry("1400x900")

        # Основные переменные
        self.model_loaded = False
        self.llm = None
        self.current_data = None
        self.current_chat_id = None
        self.chats = {}
        self.chat_messages = {}
        self.processing = False
        self.attached_file = None

        # Инициализируем алгоритмы построения графиков
        self.plot_algorithms = WellPlotAlgorithms()

        # Пути к файлам
        self.base_path = os.path.dirname(os.path.abspath(__file__))
        self.model_path = None  # Больше нет жестко заданного пути
        self.facts_path = os.path.join(self.base_path, "facts.xlsx")
        self.logo_path = os.path.join(self.base_path, "logo.png")
        self.crab_paths = {
            'thinking': os.path.join(self.base_path, "crab1.png"),
            'ready': os.path.join(self.base_path, "crab2.png"),
            'fact': os.path.join(self.base_path, "crab4.png")
        }
        self.history_path = os.path.join(self.base_path, "chats_history.json")
        self.settings_path = os.path.join(self.base_path, "settings.json")
        self.model_config_path = os.path.join(self.base_path, "model_config.json")

        # Цветовая схема
        self.colors = {
            'bg': '#1e1e2e',
            'panel': '#2b2b3b',
            'text': '#ffffff',
            'accent_blue': '#61afef',
            'accent_green': '#98c379',
            'accent_red': '#e06c75',
            'accent_yellow': '#e5c07b'
        }

        # Настройки шрифтов
        self.font_settings = {
            'family': 'Arial',
            'size': 10
        }

        # Загружаем сохраненные настройки
        self.load_settings()
        self.load_model_config()

        self.setup_ui()
        self.load_chats()

    def load_settings(self):
        """Загружает настройки интерфейса"""
        if os.path.exists(self.settings_path):
            try:
                with open(self.settings_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    self.colors.update(data.get('colors', {}))
                    self.font_settings.update(data.get('font', {}))
            except Exception as e:
                print(f"Ошибка загрузки настроек: {e}")

    def load_model_config(self):
        """Загружает путь к модели из конфига"""
        if os.path.exists(self.model_config_path):
            try:
                with open(self.model_config_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    model_path = data.get('model_path')
                    if model_path and os.path.exists(model_path):
                        self.model_path = model_path
                        print(f"Загружен путь к модели: {self.model_path}")
            except Exception as e:
                print(f"Ошибка загрузки конфига модели: {e}")

    def save_model_config(self):
        """Сохраняет путь к модели"""
        try:
            with open(self.model_config_path, 'w', encoding='utf-8') as f:
                json.dump({'model_path': self.model_path}, f, indent=2)
        except Exception as e:
            print(f"Ошибка сохранения конфига модели: {e}")

    def save_settings(self):
        """Сохраняет настройки интерфейса"""
        try:
            data = {
                'colors': self.colors,
                'font': self.font_settings
            }
            with open(self.settings_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Ошибка сохранения настроек: {e}")

    def apply_settings(self, new_colors, new_font):
        """Применяет новые настройки интерфейса"""
        self.colors = new_colors
        self.font_settings = new_font

        # Обновляем цвета всех виджетов
        self.window.configure(bg=self.colors['bg'])
        self.left_frame.configure(bg=self.colors['panel'])
        self.center_frame.configure(bg=self.colors['bg'])
        self.right_frame.configure(bg=self.colors['panel'])

        # Обновляем шрифты
        self.update_fonts()

        # Сохраняем настройки
        self.save_settings()

    def update_fonts(self):
        """Обновляет шрифты во всех виджетах"""
        # Обновляем шрифт в чате
        for widget in self.scrollable_frame.winfo_children():
            if isinstance(widget, tk.Frame):
                for child in widget.winfo_children():
                    if isinstance(child, tk.Label):
                        try:
                            child.config(
                                font=(self.font_settings['family'], self.font_settings['size'])
                            )
                        except:
                            pass

        # Обновляем поле ввода
        self.input_field.config(
            font=(self.font_settings['family'], self.font_settings['size'])
        )

    def select_model(self):
        """Открывает диалог выбора модели"""
        ModelSelectionDialog(self.window, self.model_path, self.set_model)

    def set_model(self, model_path):
        """Устанавливает новую модель и перезагружает"""
        self.model_path = model_path
        self.save_model_config()

        # Если модель уже была загружена, предлагаем перезагрузить
        if self.model_loaded:
            if messagebox.askyesno("Перезагрузка", "Модель изменена. Перезагрузить сейчас?"):
                self.reload_model()
        else:
            self.add_message("Краб",
                             f"✅ Модель выбрана: {os.path.basename(model_path)}\nНажмите '🚀 Начать работу' для загрузки")

    def reload_model(self):
        """Перезагружает модель"""
        self.model_loaded = False
        self.llm = None
        self.load_model()

    def setup_ui(self):
        """Создание трёхколоночного интерфейса"""
        self.main_paned = tk.PanedWindow(self.window, orient=tk.HORIZONTAL, bg=self.colors['bg'])
        self.main_paned.pack(fill=tk.BOTH, expand=True)

        # ===== ЛЕВАЯ ПАНЕЛЬ =====
        self.left_frame = tk.Frame(self.main_paned, bg=self.colors['panel'], width=300)
        self.main_paned.add(self.left_frame, width=300)

        self.logo_label = tk.Label(self.left_frame, bg=self.colors['panel'])
        self.load_logo()
        self.logo_label.pack(pady=20)

        self.new_chat_btn = tk.Button(self.left_frame, text="➕ New Chat",
                                      command=self.create_new_chat,
                                      bg=self.colors['accent_blue'], fg='white',
                                      font=("Arial", 12), pady=8, borderwidth=0)
        self.new_chat_btn.pack(pady=10, padx=10, fill=tk.X)

        self.chats_frame = tk.Frame(self.left_frame, bg=self.colors['panel'])
        self.chats_frame.pack(fill=tk.BOTH, expand=True, padx=5, pady=5)

        self.chat_buttons = []
        self.update_chat_list()

        # ===== ЦЕНТРАЛЬНАЯ ПАНЕЛЬ =====
        self.center_frame = tk.Frame(self.main_paned, bg=self.colors['bg'])
        self.main_paned.add(self.center_frame, width=700)

        self.setup_chat_area()

        # ===== ПРАВАЯ ПАНЕЛЬ =====
        self.right_frame = tk.Frame(self.main_paned, bg=self.colors['panel'], width=400)
        self.main_paned.add(self.right_frame, width=400)

        self.setup_crab_area()

        self.show_start_screen()

    def load_logo(self):
        """Загрузка логотипа"""
        try:
            if os.path.exists(self.logo_path):
                img = Image.open(self.logo_path)
                img = img.resize((200, 80), Image.Resampling.LANCZOS)
                self.logo_img = ImageTk.PhotoImage(img)
                self.logo_label.config(image=self.logo_img)
            else:
                self.logo_label.config(text="🦀 НЕФТЯНОЙ\n   КРАБ", font=("Arial", 16, "bold"),
                                       fg=self.colors['accent_yellow'], bg=self.colors['panel'])
        except:
            self.logo_label.config(text="🦀 НЕФТЯНОЙ\n   КРАБ", font=("Arial", 16, "bold"),
                                   fg=self.colors['accent_yellow'], bg=self.colors['panel'])

    def setup_chat_area(self):
        """Настройка области чата"""
        self.chat_title = tk.Label(self.center_frame, text="Чат",
                                   font=("Arial", 14, "bold"),
                                   bg=self.colors['bg'], fg=self.colors['text'])
        self.chat_title.pack(pady=10)

        self.messages_frame = tk.Frame(self.center_frame, bg=self.colors['bg'])
        self.messages_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)

        self.messages_canvas = tk.Canvas(self.messages_frame, bg=self.colors['bg'],
                                         highlightthickness=0)
        self.messages_scrollbar = tk.Scrollbar(self.messages_frame, orient=tk.VERTICAL,
                                               command=self.messages_canvas.yview)
        self.scrollable_frame = tk.Frame(self.messages_canvas, bg=self.colors['bg'])

        self.scrollable_frame.bind("<Configure>",
                                   lambda e: self.messages_canvas.configure(
                                       scrollregion=self.messages_canvas.bbox("all")))

        self.messages_canvas.create_window((0, 0), window=self.scrollable_frame, anchor="nw")
        self.messages_canvas.configure(yscrollcommand=self.messages_scrollbar.set)

        self.messages_canvas.pack(side=tk.LEFT, fill=tk.BOTH, expand=True)
        self.messages_scrollbar.pack(side=tk.RIGHT, fill=tk.Y)

        # Нижняя панель ввода с кнопкой прикрепления
        input_frame = tk.Frame(self.center_frame, bg=self.colors['bg'])
        input_frame.pack(fill=tk.X, padx=10, pady=10)

        attach_btn = tk.Button(input_frame, text="📎",
                               command=self.attach_file,
                               bg=self.colors['panel'], fg=self.colors['text'],
                               font=("Arial", 12), padx=10, borderwidth=0)
        attach_btn.pack(side=tk.LEFT, padx=(0, 5))

        self.attached_label = tk.Label(input_frame, text="",
                                       bg=self.colors['bg'], fg=self.colors['accent_green'])
        self.attached_label.pack(side=tk.LEFT, padx=5)

        self.input_field = tk.Text(input_frame, height=3,
                                   font=(self.font_settings['family'], self.font_settings['size']),
                                   bg=self.colors['panel'],
                                   fg=self.colors['text'],
                                   insertbackground=self.colors['text'],
                                   wrap=tk.WORD)
        self.input_field.pack(side=tk.LEFT, fill=tk.X, expand=True, padx=(0, 10))
        self.input_field.bind('<Return>', self.on_enter_pressed)

        send_btn = tk.Button(input_frame, text="➤ Отправить",
                             command=self.send_message,
                             bg=self.colors['accent_blue'], fg='white',
                             font=("Arial", 11), padx=15, pady=5, borderwidth=0)
        send_btn.pack(side=tk.RIGHT)

    def on_enter_pressed(self, event):
        if not event.state & 0x1:
            self.send_message()
            return "break"

    def setup_crab_area(self):
        """Настройка области краба"""
        # Верхняя часть с крабом
        crab_top_frame = tk.Frame(self.right_frame, bg=self.colors['panel'])
        crab_top_frame.pack(fill='x', padx=10, pady=5)

        tk.Label(crab_top_frame, text="🦀 Краб-помощник",
                 font=("Arial", 14, "bold"),
                 bg=self.colors['panel'], fg=self.colors['accent_yellow']).pack(side='left')

        # Кнопка выбора модели
        model_btn = tk.Button(
            crab_top_frame,
            text="📦 Модель",
            command=self.select_model,
            bg=self.colors['accent_blue'],
            fg='white',
            font=("Arial", 10),
            padx=10,
            borderwidth=0,
            cursor="hand2"
        )
        model_btn.pack(side='right', padx=5)

        # Кнопка настроек
        settings_btn = tk.Button(
            crab_top_frame,
            text="⚙️",
            command=self.open_settings,
            bg=self.colors['panel'],
            fg=self.colors['text'],
            font=("Arial", 12),
            borderwidth=0,
            cursor="hand2"
        )
        settings_btn.pack(side='right', padx=5)

        self.crab_frame = tk.Frame(self.right_frame, bg=self.colors['panel'],
                                   width=350, height=350)
        self.crab_frame.pack(pady=10)
        self.crab_frame.pack_propagate(False)

        self.crab_label = tk.Label(self.crab_frame, bg=self.colors['panel'])
        self.crab_label.pack(expand=True)

        self.load_crab_image('ready')

        self.fact_btn = tk.Button(self.right_frame, text="🎲 Интересно",
                                  command=self.show_random_fact,
                                  bg=self.colors['accent_green'], fg='white',
                                  font=("Arial", 14), pady=8, borderwidth=0,
                                  state=tk.NORMAL)
        self.fact_btn.pack(pady=10, padx=20, fill=tk.X)

        self.fact_text = tk.Text(self.right_frame, height=6,
                                 font=("Arial", 14),
                                 bg=self.colors['panel'],
                                 fg=self.colors['text'],
                                 wrap=tk.WORD, borderwidth=0)
        self.fact_text.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        self.fact_text.insert(1.0, "Нажми 'Интересно' для случайного факта о нефти!")
        self.fact_text.config(state=tk.DISABLED)

        # Индикатор модели
        self.model_indicator = tk.Label(
            self.right_frame,
            text="⚪ Модель не выбрана",
            bg=self.colors['panel'],
            fg=self.colors['accent_red'],
            font=("Arial", 9)
        )
        self.model_indicator.pack(pady=5)

    def update_model_indicator(self):
        """Обновляет индикатор состояния модели"""
        if self.model_loaded:
            self.model_indicator.config(
                text=f"✅ Модель загружена: {os.path.basename(self.model_path)[:30]}...",
                fg=self.colors['accent_green']
            )
        elif self.model_path:
            self.model_indicator.config(
                text=f"⏳ Модель выбрана: {os.path.basename(self.model_path)[:30]}...",
                fg=self.colors['accent_yellow']
            )
        else:
            self.model_indicator.config(
                text="⚪ Модель не выбрана",
                fg=self.colors['accent_red']
            )

    def open_settings(self):
        """Открывает окно настроек"""
        SettingsDialog(self.window, self.colors, self.font_settings, self.apply_settings)

    def load_crab_image(self, state):
        """Загрузка изображения краба"""
        try:
            if state == 'thinking':
                path = self.crab_paths['thinking']
            elif state == 'ready':
                path = self.crab_paths['ready']
            elif state == 'fact':
                path = self.crab_paths['fact']
            else:
                return

            if os.path.exists(path):
                img = Image.open(path)
                img.thumbnail((300, 300), Image.Resampling.LANCZOS)
                self.crab_img = ImageTk.PhotoImage(img)
                self.crab_label.config(image=self.crab_img)
                return True
            else:
                self.crab_label.config(text=f"🦀", font=("Arial", 80),
                                       fg=self.colors['accent_yellow'])
                return False
        except Exception as e:
            print(f"Ошибка загрузки изображения краба: {e}")
            self.crab_label.config(text=f"🦀", font=("Arial", 80),
                                   fg=self.colors['accent_yellow'])
            return False

    def update_crab_image(self, state):
        self.window.after(0, lambda: self._update_crab_image(state))

    def _update_crab_image(self, state):
        self.load_crab_image(state)

    def show_start_screen(self):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()

        # Проверяем, выбрана ли модель
        if not self.model_path:
            text = "⚠️ Сначала выберите файл модели (кнопка '📦 Модель')"
            btn_text = "📦 Выбрать модель"
            btn_cmd = self.select_model
            btn_color = self.colors['accent_yellow']
        else:
            text = "🚀 Модель выбрана! Можно начинать работу"
            btn_text = "🚀 Начать работу"
            btn_cmd = self.start_model_and_chat
            btn_color = self.colors['accent_green']

        info_label = tk.Label(
            self.scrollable_frame,
            text=text,
            font=("Arial", 14),
            bg=self.colors['bg'],
            fg=self.colors['text']
        )
        info_label.pack(pady=50)

        start_btn = tk.Button(self.scrollable_frame,
                              text=btn_text,
                              command=btn_cmd,
                              bg=btn_color,
                              fg='white',
                              font=("Arial", 18, "bold"),
                              padx=50, pady=20,
                              borderwidth=0,
                              cursor="hand2")
        start_btn.pack(expand=True, pady=50)

        self.update_model_indicator()

    def start_model_and_chat(self):
        if not self.model_path:
            messagebox.showwarning("Внимание", "Сначала выберите файл модели!")
            self.select_model()
            return
        self.create_new_chat()
        self.load_model()

    def load_model(self):
        if not os.path.exists(self.model_path):
            messagebox.showerror("Ошибка", f"Модель не найдена:\n{self.model_path}")
            return

        def load():
            self.processing = True
            self.update_crab_image('thinking')
            self.fact_btn.config(state=tk.DISABLED)

            try:
                self.add_message("Краб",
                                 f"⏳ Загружаю модель {os.path.basename(self.model_path)}... Это может занять несколько минут")

                self.llm = Llama(
                    model_path=self.model_path,
                    n_ctx=2048,
                    n_threads=4,
                    n_gpu_layers=0,
                    verbose=False
                )
                self.model_loaded = True
                self.add_message("Краб", "✅ Модель загружена! Задавайте вопросы.")
                self.update_model_indicator()
            except Exception as e:
                messagebox.showerror("Ошибка", f"Не удалось загрузить модель:\n{str(e)}")
                self.add_message("Краб", f"❌ Ошибка загрузки модели: {str(e)}")
            finally:
                self.processing = False
                self.update_crab_image('ready')
                self.fact_btn.config(state=tk.NORMAL)

        thread = threading.Thread(target=load)
        thread.daemon = True
        thread.start()

    def attach_file(self):
        """Прикрепление Excel файла с данными"""
        file_path = filedialog.askopenfilename(
            title="Выберите Excel файл с данными",
            filetypes=[("Excel files", "*.xlsx *.xls"), ("CSV files", "*.csv"), ("All files", "*.*")]
        )
        if file_path:
            try:
                if file_path.endswith('.csv'):
                    self.current_data = pd.read_csv(file_path)
                else:
                    self.current_data = pd.read_excel(file_path)

                # Автоопределение колонок
                self.plot_algorithms.auto_detect_columns(self.current_data)
                self.current_data = self.plot_algorithms.prepare_data(self.current_data)

                self.attached_file = file_path
                self.attached_label.config(text=f"📊 {os.path.basename(file_path)}")

                # Показываем информацию о файле
                wells = self.current_data[self.plot_algorithms.column_names['well']].nunique()
                self.add_message("Краб", f"✅ Файл загружен: {os.path.basename(file_path)}\n"
                                         f"   Записей: {len(self.current_data)}, "
                                         f"Скважин: {wells}")

                # Показываем первые 5 скважин
                sample_wells = self.current_data[self.plot_algorithms.column_names['well']].unique()[:5]
                self.add_message("Краб", f"📋 Примеры скважин: {', '.join(map(str, sample_wells))}")

            except Exception as e:
                messagebox.showerror("Ошибка", f"Не удалось загрузить файл:\n{str(e)}")
                traceback.print_exc()

    def show_random_fact(self):
        if self.processing:
            return

        try:
            if os.path.exists(self.facts_path):
                df = pd.read_excel(self.facts_path)
                if len(df.columns) >= 2:
                    facts = df.iloc[:, 1].tolist()
                    if facts:
                        fact = random.choice(facts)
                        self.update_crab_image('fact')
                        self.fact_text.config(state=tk.NORMAL)
                        self.fact_text.delete(1.0, tk.END)
                        self.fact_text.insert(1.0, f"А ты знал, что...\n\n{fact}")
                        self.fact_text.config(state=tk.DISABLED)
        except Exception as e:
            print(f"Ошибка загрузки фактов: {e}")

    def create_new_chat(self):
        chat_id = len(self.chats)
        chat_name = f"Чат {chat_id + 1}"
        self.chats[chat_id] = chat_name
        self.chat_messages[chat_id] = []
        self.current_chat_id = chat_id
        self.update_chat_list()
        self.clear_chat_display()
        self.save_chats()

    def update_chat_list(self):
        for btn in self.chat_buttons:
            btn.destroy()
        self.chat_buttons.clear()

        for chat_id, chat_name in self.chats.items():
            frame = tk.Frame(self.chats_frame, bg=self.colors['panel'])
            frame.pack(fill=tk.X, pady=2)

            btn = tk.Button(frame, text=f"  {chat_name}",
                            command=lambda cid=chat_id: self.switch_chat(cid),
                            bg=self.colors['panel'],
                            fg=self.colors['text'],
                            font=("Arial", 10),
                            anchor=tk.W, borderwidth=0)
            btn.pack(side=tk.LEFT, fill=tk.X, expand=True)

            btn.bind('<Double-Button-1>', lambda e, cid=chat_id: self.rename_chat(cid))

            self.chat_buttons.append(frame)

    def switch_chat(self, chat_id):
        self.current_chat_id = chat_id
        self.display_chat_messages()

    def rename_chat(self, chat_id):
        dialog = tk.Toplevel(self.window)
        dialog.title("Переименовать чат")
        dialog.geometry("300x100")
        dialog.configure(bg=self.colors['panel'])

        tk.Label(dialog, text="Новое название:",
                 bg=self.colors['panel'], fg=self.colors['text']).pack(pady=10)

        entry = tk.Entry(dialog, width=30)
        entry.pack(pady=5)
        entry.insert(0, self.chats[chat_id])

        def save():
            new_name = entry.get().strip()
            if new_name:
                self.chats[chat_id] = new_name
                self.update_chat_list()
                self.save_chats()
            dialog.destroy()

        tk.Button(dialog, text="Сохранить", command=save,
                  bg=self.colors['accent_blue'], fg='white').pack(pady=5)

    def clear_chat_display(self):
        for widget in self.scrollable_frame.winfo_children():
            widget.destroy()

    def display_chat_messages(self):
        self.clear_chat_display()
        messages = self.chat_messages.get(self.current_chat_id, [])
        for msg in messages:
            self.display_message(msg['sender'], msg['text'], msg.get('image'))

    def display_message(self, sender, text, image_path=None):
        msg_frame = tk.Frame(self.scrollable_frame, bg=self.colors['bg'])
        msg_frame.pack(fill=tk.X, pady=5)

        if sender == "Вы":
            label = tk.Label(msg_frame, text=f"👤 {sender}:",
                             bg=self.colors['bg'], fg=self.colors['accent_green'],
                             font=(self.font_settings['family'], self.font_settings['size'], "bold"))
        else:
            label = tk.Label(msg_frame, text=f"🦀 {sender}:",
                             bg=self.colors['bg'], fg=self.colors['accent_yellow'],
                             font=(self.font_settings['family'], self.font_settings['size'], "bold"))
        label.pack(anchor=tk.W)

        if image_path and os.path.exists(image_path):
            try:
                img = Image.open(image_path)
                img.thumbnail((500, 300), Image.Resampling.LANCZOS)
                photo = ImageTk.PhotoImage(img)

                img_label = tk.Label(msg_frame, image=photo, bg=self.colors['bg'])
                img_label.image = photo
                img_label.pack(pady=5)

                img_label.bind('<Button-1>', lambda e, p=image_path: self.show_full_image(p))
                img_label.bind('<Button-3>', lambda e, p=image_path: self.copy_image(p))

            except Exception as e:
                tk.Label(msg_frame, text=f"[Ошибка загрузки изображения: {e}]",
                         bg=self.colors['bg'], fg=self.colors['accent_red']).pack()

        text_label = tk.Label(msg_frame, text=text,
                              bg=self.colors['bg'], fg=self.colors['text'],
                              font=(self.font_settings['family'], self.font_settings['size']),
                              wraplength=600, justify=tk.LEFT)
        text_label.pack(anchor=tk.W)
        text_label.bind('<Button-3>', lambda e, t=text: self.copy_text(t))

    def show_full_image(self, image_path):
        try:
            img = Image.open(image_path)
            img.show()
        except Exception as e:
            messagebox.showerror("Ошибка", f"Не удалось открыть изображение:\n{str(e)}")

    def copy_image(self, image_path):
        try:
            import subprocess
            subprocess.run(
                ['osascript', '-e', f'set the clipboard to (read (POSIX file "{image_path}") as JPEG picture)'])
            messagebox.showinfo("Успех", "Изображение скопировано в буфер обмена")
        except:
            messagebox.showwarning("Внимание", "Копирование изображения не поддерживается на этой системе")

    def copy_text(self, text):
        self.window.clipboard_clear()
        self.window.clipboard_append(text)
        messagebox.showinfo("Успех", "Текст скопирован в буфер обмена")

    def add_message(self, sender, text, image_path=None):
        if self.current_chat_id is None:
            self.create_new_chat()

        message = {'sender': sender, 'text': text, 'image': image_path}
        self.chat_messages.setdefault(self.current_chat_id, []).append(message)
        self.display_message(sender, text, image_path)
        self.save_chats()

    def extract_well_name(self, text):
        """Извлечение номера скважины из текста"""
        text = text.lower().strip()

        print(f"DEBUG: Извлекаем номер скважины из: '{text}'")

        # Паттерны для поиска
        patterns = [
            r'скв[ажин]?[а-я]*[-\s]*(\d+)',
            r'well[-\s]*(\d+)',
            r'(\d+)[\s]*(?:скв|well)',
            r'^(\d{3,4})$',
            r'(\d{3,4})$',
            r'(?:для|по)[\s]*(\d+)',
            r'(\d{3,4})'
        ]

        for i, pattern in enumerate(patterns):
            match = re.search(pattern, text)
            if match:
                well_num = match.group(1)
                result = f"СКВ-{well_num}"
                print(f"DEBUG: Паттерн {i} сработал: {well_num} -> {result}")
                return result

        print("DEBUG: Номер скважины не найден")
        return None

    def send_message(self):
        if self.current_chat_id is None:
            messagebox.showwarning("Внимание", "Сначала создайте новый чат!")
            return

        user_text = self.input_field.get(1.0, tk.END).strip()
        if not user_text:
            return

        self.input_field.delete(1.0, tk.END)
        self.add_message("Вы", user_text)

        # Проверяем, команда ли это для построения графика
        if self.current_data is not None and self.is_plot_command(user_text):
            thread = threading.Thread(target=self.generate_plot, args=(user_text,))
            thread.daemon = True
            thread.start()
        # Проверяем, команда ли это для создания презентации
        elif self.is_presentation_command(user_text):
            thread = threading.Thread(target=self.generate_presentation, args=(user_text,))
            thread.daemon = True
            thread.start()
        else:
            # Обычный вопрос к ИИ
            if not self.model_loaded:
                self.add_message("Краб", "❌ Модель не загружена. Сначала выберите и загрузите модель!")
                return
            thread = threading.Thread(target=self.generate_response, args=(user_text,))
            thread.daemon = True
            thread.start()

    def is_plot_command(self, text):
        """Проверка, является ли текст командой для построения графика"""
        cmd_lower = text.lower()
        plot_keywords = ['график', 'построй', 'нарисуй', 'покажи',
                         'чен', 'ipr', 'decline', 'падени', 'дебит', 'обводн']
        result = any(word in cmd_lower for word in plot_keywords)
        print(f"DEBUG: is_plot_command({text}) = {result}")
        return result

    def is_presentation_command(self, text):
        """Проверка, является ли текст командой для создания презентации"""
        cmd_lower = text.lower()
        presentation_keywords = ['презентац', 'сделай отчет', 'сформируй отчет',
                                 'все графики', 'полный отчет', 'слайд']
        result = any(word in cmd_lower for word in presentation_keywords)
        print(f"DEBUG: is_presentation_command({text}) = {result}")
        return result

    def generate_plot(self, command):
        """Генерация графика по команде"""
        self.processing = True
        self.update_crab_image('thinking')
        self.fact_btn.config(state=tk.DISABLED)

        try:
            cmd_lower = command.lower()
            well_name = self.extract_well_name(command)

            print(f"DEBUG: generate_plot: command='{command}', well_name='{well_name}'")

            if not well_name:
                self.add_message("Краб", "❌ Укажите номер скважины, например: 'чен 109'")
                return

            # Проверяем наличие данных
            if self.current_data is None:
                self.add_message("Краб", "❌ Сначала загрузите файл с данными (кнопка 📎)")
                return

            well_col = self.plot_algorithms.column_names['well']
            available_wells = self.current_data[well_col].unique()

            print(f"DEBUG: Доступные скважины: {available_wells}")

            if well_name not in available_wells:
                self.add_message("Краб", f"❌ Скважина {well_name} не найдена в загруженных данных")
                self.add_message("Краб", f"📋 Доступные скважины: {', '.join(map(str, available_wells[:10]))}")
                return

            # Создаем временный файл для графика
            temp_file = os.path.join(tempfile.gettempdir(), f"crab_plot_{random.randint(1000, 9999)}.png")

            # Определяем тип графика
            if any(word in cmd_lower for word in ['чен', 'chen']):
                img_path, analysis = self.plot_algorithms.plot_chen(self.current_data, well_name, temp_file)
                self.add_message("Краб", f"📊 График Чена для {well_name}", img_path)
                self.add_message("Краб", analysis)

            elif any(word in cmd_lower for word in ['ipr', 'индикаторн']):
                img_path, analysis = self.plot_algorithms.plot_ipr(self.current_data, well_name, temp_file)
                self.add_message("Краб", f"📈 IPR диаграмма для {well_name}", img_path)

            elif any(word in cmd_lower for word in ['decline', 'падени']):
                img_path, analysis = self.plot_algorithms.plot_decline(self.current_data, well_name, temp_file)
                self.add_message("Краб", f"📉 Кривая падения для {well_name}", img_path)

            elif any(word in cmd_lower for word in ['дебит', 'нефт']):
                img_path, analysis = self.plot_algorithms.plot_standard(self.current_data, well_name, temp_file,
                                                                        'oil_rate')
                self.add_message("Краб", f"📊 Дебит нефти: {well_name}", img_path)

            elif any(word in cmd_lower for word in ['обводн', 'water']):
                img_path, analysis = self.plot_algorithms.plot_standard(self.current_data, well_name, temp_file,
                                                                        'watercut')
                self.add_message("Краб", f"📊 Обводненность: {well_name}", img_path)

            else:
                self.add_message("Краб", "❌ Неизвестный тип графика. Доступны: чен, ipr, decline, дебит, обводненность")

        except Exception as e:
            self.add_message("Краб", f"❌ Ошибка построения графика: {str(e)}")
            traceback.print_exc()

        finally:
            self.processing = False
            self.update_crab_image('ready')
            self.fact_btn.config(state=tk.NORMAL)

    def generate_presentation(self, command):
        """Генерация презентации со всеми графиками по скважине"""
        self.processing = True
        self.update_crab_image('thinking')
        self.fact_btn.config(state=tk.DISABLED)

        try:
            well_name = self.extract_well_name(command)

            if not well_name:
                self.add_message("Краб", "❌ Укажите номер скважины, например: 'сделай презентацию по 109'")
                return

            # Проверяем наличие данных
            if self.current_data is None:
                self.add_message("Краб", "❌ Сначала загрузите файл с данными (кнопка 📎)")
                return

            well_col = self.plot_algorithms.column_names['well']
            available_wells = self.current_data[well_col].unique()

            if well_name not in available_wells:
                self.add_message("Краб", f"❌ Скважина {well_name} не найдена в загруженных данных")
                return

            self.add_message("Краб", f"🔄 Создаю презентацию для {well_name}... Это займет некоторое время")

            # Создаем временную папку для графиков
            with tempfile.TemporaryDirectory() as tmpdir:
                chart_paths = []
                analyses = {}

                # 1. График Чена
                chen_file = os.path.join(tmpdir, "01_chen.png")
                chen_path, chen_analysis = self.plot_algorithms.plot_chen(self.current_data, well_name, chen_file)
                chart_paths.append(("График Чена", chen_path))
                analyses["chen"] = chen_analysis

                # 2. IPR диаграмма
                ipr_file = os.path.join(tmpdir, "02_ipr.png")
                ipr_path, ipr_analysis = self.plot_algorithms.plot_ipr(self.current_data, well_name, ipr_file)
                chart_paths.append(("IPR диаграмма", ipr_path))
                analyses["ipr"] = ipr_analysis

                # 3. Decline curve
                decline_file = os.path.join(tmpdir, "03_decline.png")
                decline_path, decline_analysis = self.plot_algorithms.plot_decline(self.current_data, well_name,
                                                                                   decline_file)
                chart_paths.append(("Кривая падения", decline_path))
                analyses["decline"] = decline_analysis

                # 4. Дебит нефти
                oil_file = os.path.join(tmpdir, "04_oil_rate.png")
                oil_path, oil_analysis = self.plot_algorithms.plot_standard(self.current_data, well_name, oil_file,
                                                                            'oil_rate')
                chart_paths.append(("Дебит нефти", oil_path))
                analyses["oil"] = oil_analysis

                # 5. Обводненность
                water_file = os.path.join(tmpdir, "05_watercut.png")
                water_path, water_analysis = self.plot_algorithms.plot_standard(self.current_data, well_name,
                                                                                water_file, 'watercut')
                chart_paths.append(("Обводненность", water_path))
                analyses["water"] = water_analysis

                # Создаем презентацию
                self.add_message("Краб", "📊 Собираю слайды...")
                pptx_path = self.create_pptx(well_name, chart_paths, analyses, tmpdir)

                # Спрашиваем, куда сохранить
                save_path = filedialog.asksaveasfilename(
                    title="Сохранить презентацию",
                    defaultextension=".pptx",
                    filetypes=[("PowerPoint", "*.pptx")],
                    initialfile=f"{well_name}_отчет_{datetime.now().strftime('%Y%m%d')}.pptx"
                )

                if save_path:
                    shutil.copy2(pptx_path, save_path)
                    self.add_message("Краб", f"✅ Презентация сохранена:\n{save_path}")

                    # Показываем preview в чате
                    self.add_message("Краб", f"📁 Файл готов: {os.path.basename(save_path)}")
                else:
                    self.add_message("Краб", "❌ Сохранение отменено")

        except Exception as e:
            self.add_message("Краб", f"❌ Ошибка создания презентации: {str(e)}")
            traceback.print_exc()

        finally:
            self.processing = False
            self.update_crab_image('ready')
            self.fact_btn.config(state=tk.NORMAL)

    def create_pptx(self, well_name, chart_paths, analyses, tmpdir):
        """Создает PPTX файл с графиками и анализом"""
        prs = Presentation()

        # Титульный слайд
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = f"Анализ скважины {well_name}"
        subtitle.text = f"Дата отчета: {datetime.now().strftime('%d.%m.%Y')}\n"
        subtitle.text += f"Выполнил: 🦀 Нефтяной ассистент"

        # Слайд с общей информацией
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        content = slide.placeholders[1]

        title.text = "Общая информация"

        well_data = self.current_data[self.current_data[self.plot_algorithms.column_names['well']] == well_name]
        oil_rates = well_data[self.plot_algorithms.column_names['oil_rate']].values
        watercuts = well_data[self.plot_algorithms.column_names['watercut']].values if \
        self.plot_algorithms.column_names['watercut'] in well_data.columns else None

        info_text = f"Скважина: {well_name}\n"
        info_text += f"Количество замеров: {len(well_data)}\n"
        info_text += f"Период: с {well_data[self.plot_algorithms.column_names['date']].min()} по {well_data[self.plot_algorithms.column_names['date']].max()}\n\n"
        info_text += f"Средний дебит нефти: {np.mean(oil_rates):.1f} т/сут\n"
        if watercuts is not None:
            info_text += f"Средняя обводненность: {np.mean(watercuts):.1f}%\n"

        content.text = info_text

        # Слайды с графиками
        for title_text, chart_path in chart_paths:
            slide_layout = prs.slide_layouts[6]  # Пустой слайд
            slide = prs.slides.add_slide(slide_layout)

            # Заголовок
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.8))
            tf = txBox.text_frame
            tf.text = f"{title_text} - {well_name}"
            tf.paragraphs[0].font.size = Inches(0.3)
            tf.paragraphs[0].font.bold = True

            # График
            left = Inches(0.5)
            top = Inches(1.0)
            width = Inches(9)
            height = Inches(5)

            slide.shapes.add_picture(chart_path, left, top, width=width, height=height)

            # Для графика Чена добавляем анализ
            if title_text == "График Чена" and "chen" in analyses:
                analysis_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(6.2), Inches(9), Inches(1.5)
                )
                tf = analysis_box.text_frame
                tf.text = analyses["chen"][:500] + "..."  # Обрезаем для слайда
                tf.paragraphs[0].font.size = Inches(0.15)

        # Слайд с детальным анализом Чена
        if "chen" in analyses:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            title = slide.shapes.title
            content = slide.placeholders[1]

            title.text = "Детальный анализ обводнения"
            content.text = analyses["chen"]

        # Сохраняем
        pptx_path = os.path.join(tmpdir, f"{well_name}_report.pptx")
        prs.save(pptx_path)

        return pptx_path

    def generate_response(self, question):
        """Обычный текстовый ответ от ИИ"""
        self.processing = True
        self.update_crab_image('thinking')
        self.fact_btn.config(state=tk.DISABLED)

        try:
            prompt = f"Вопрос по нефтегазовой геологии: {question}\n\nОтвет:"

            output = self.llm(
                prompt,
                max_tokens=300,
                temperature=0.7,
                stop=["\n\n", "Вопрос:"],
                echo=False
            )

            response = output["choices"][0]["text"].strip()
            if not response:
                response = "Извини, не могу ответить на этот вопрос."

            self.add_message("Краб", response)

        except Exception as e:
            self.add_message("Краб", f"❌ Ошибка: {str(e)[:200]}")

        finally:
            self.processing = False
            self.update_crab_image('ready')
            self.fact_btn.config(state=tk.NORMAL)

    def load_chats(self):
        if os.path.exists(self.history_path):
            try:
                with open(self.history_path, 'r', encoding='utf-8') as f:
                    data = json.load(f)
                    self.chats = {int(k): v for k, v in data.get('chats', {}).items()}
                    self.chat_messages = {int(k): v for k, v in data.get('messages', {}).items()}
                    if self.chats:
                        self.current_chat_id = max(self.chats.keys())
            except Exception as e:
                print(f"Ошибка загрузки чатов: {e}")

    def save_chats(self):
        try:
            data = {
                'chats': self.chats,
                'messages': self.chat_messages
            }
            with open(self.history_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, ensure_ascii=False, indent=2)
        except Exception as e:
            print(f"Ошибка сохранения чатов: {e}")

    def run(self):
        self.window.mainloop()


if __name__ == "__main__":
    app = CrabAssistant()
    app.run()
